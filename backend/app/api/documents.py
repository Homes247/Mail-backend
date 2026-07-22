import json
import uuid
import os
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from typing import Optional
import asyncio

from fastapi import APIRouter, Depends, HTTPException
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, or_
from pydantic import BaseModel

from app.database import get_db
from app.models.document import Document
from app.models.document_share import DocumentShare
from app.models.user import User
from app.models.audit_event import AuditEvent
from app.models.sheet_version import SheetVersion
from app.api.auth import get_current_user, get_optional_current_user
from app.lib.document_storage import DocumentStorage
from fastapi import UploadFile, File, Form
import datetime as _dt
import io

router = APIRouter()


class DocumentCreate(BaseModel):
    title: str = "Untitled"
    doc_type: str  # sheet | doc | slide


class DocumentUpdate(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None
    is_trashed: Optional[int] = None

class DocumentShareCreate(BaseModel):
    email: str
    permission: str


class VersionResponse(BaseModel):
    id: str
    version_name: Optional[str] = None
    created_at: _dt.datetime
    created_by_user_id: int
    is_named: bool
    contributors: list

class VersionCreateRequest(BaseModel):
    version_name: str


def _default_content(doc_type: str, title: str) -> str:
    if doc_type == "doc":
        return json.dumps({"html": "<br>"})
    if doc_type == "sheet":
        return json.dumps({"cells": {}, "formats": {"0,0": {"size": "13px"}}})
    if doc_type == "slide":
        page_id = str(uuid.uuid4())
        return json.dumps({
            "id":        page_id,
            "title":     title,
            "pages":     {page_id: {"title": "", "body": ""}},
            "pageOrder": [page_id],
        })
    return "{}"


@router.post("")
async def create_document(
    body: DocumentCreate,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    content_val = _default_content(body.doc_type, body.title)
    doc = Document(
        title=body.title,
        doc_type=body.doc_type,
        owner_id=current_user.id,
    )
    db.add(doc)
    await db.commit()
    await db.refresh(doc)
    
    try:
        storage_res = await asyncio.to_thread(
            DocumentStorage.save, doc.owner_id, doc.id, content_val, doc_type=doc.doc_type
        )
        doc.file_path = storage_res["relative_path"]
        doc.file_size = storage_res["size"]
        doc.content_version = 1
        await db.commit()
    except Exception as e:
        import traceback
        print(f"Error saving document to storage: {e}")
        print(traceback.format_exc())
        
    return {"id": doc.id, "title": doc.title, "doc_type": doc.doc_type}

def resolve_validation_options(formula, wb, ws):
    try:
        with open("backend/debug_validation.log", "a") as f_log:
            f_log.write(f"Resolving validation: formula={formula}, sheet={ws.title}\n")
    except:
        pass
        
    if not formula:
        return []
    formula = formula.strip('"')
    if not formula.startswith('='):
        # Comma-separated list
        return [{"label": x.strip()} for x in formula.split(',') if x.strip()]
    
    # It is a reference. E.g. =Sheet2!$A$1:$A$10 or =$A$1:$A$10
    ref = formula[1:] # strip '='
    sheet_name = ws.title
    if '!' in ref:
        sheet_name, ref = ref.split('!', 1)
        sheet_name = sheet_name.strip("'")
    
    ref = ref.replace('$', '') # remove absolute references
    
    try:
        with open("backend/debug_validation.log", "a") as f_log:
            f_log.write(f"  Parsed sheet_name={sheet_name}, ref={ref}\n")
    except:
        pass

    if sheet_name in wb.sheetnames:
        target_ws = wb[sheet_name]
        try:
            cells_range = target_ws[ref]
            options = []
            if hasattr(cells_range, 'value'): # Single cell
                val = cells_range.value
                if val is not None:
                    options.append({"label": str(val)})
            else:
                for row in cells_range:
                    if isinstance(row, tuple) or isinstance(row, list):
                        for cell in row:
                            if cell.value is not None:
                                options.append({"label": str(cell.value)})
                    elif hasattr(row, 'value') and row.value is not None:
                        options.append({"label": str(row.value)})
            
            try:
                with open("backend/debug_validation.log", "a") as f_log:
                    f_log.write(f"  Success, options={options}\n")
            except:
                pass
            return options
        except Exception as e:
            try:
                with open("backend/debug_validation.log", "a") as f_log:
                    f_log.write(f"  Error resolving range: {e}\n")
            except:
                pass
            return []
    else:
        try:
            with open("backend/debug_validation.log", "a") as f_log:
                f_log.write(f"  Sheet name {sheet_name} not found in {wb.sheetnames}\n")
        except:
            pass
    return []

def repair_zoho_rich_values(content_bytes: bytes) -> bytes:
    import zipfile
    import xml.etree.ElementTree as ET
    import os
    import tempfile
    import shutil
    import openpyxl
    from openpyxl.drawing.image import Image as OpenpyxlImage
    
    NS = {
        'main': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'xlrd': 'http://schemas.microsoft.com/office/spreadsheetml/2017/richdata',
        'xldr': 'http://schemas.microsoft.com/office/spreadsheetml/2017/richdata2',
        'xlrr': 'http://schemas.microsoft.com/office/spreadsheetml/2022/richvaluerel'
    }
    
    temp_dir = tempfile.mkdtemp()
    input_path = os.path.join(temp_dir, "input.xlsx")
    output_path = os.path.join(temp_dir, "output.xlsx")
    
    try:
        with open(input_path, 'wb') as f:
            f.write(content_bytes)
            
        with zipfile.ZipFile(input_path, 'r') as z:
            namelist = z.namelist()
            sheet_targets = {}
            if 'xl/_rels/workbook.xml.rels' in namelist:
                root = ET.fromstring(z.read('xl/_rels/workbook.xml.rels'))
                for rel in root.findall('.//r:Relationship', namespaces={'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}):
                    sheet_targets[rel.attrib['Id']] = 'xl/' + rel.attrib['Target']
            sheet_info = {}
            if 'xl/workbook.xml' in namelist:
                root = ET.fromstring(z.read('xl/workbook.xml'))
                for sheet in root.findall('.//main:sheet', namespaces=NS):
                    sheet_info[sheet.attrib['name']] = sheet_targets.get(sheet.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'])
            rc_list = []
            future_metadata_bk = []
            if 'xl/metadata.xml' in namelist:
                root = ET.fromstring(z.read('xl/metadata.xml'))
                value_meta = root.find('.//main:valueMetadata', namespaces=NS)
                if value_meta is not None:
                    for bk in value_meta.findall('.//main:bk', namespaces=NS):
                        for rc in bk.findall('.//main:rc', namespaces=NS):
                            rc_list.append(int(rc.attrib['v']))
                future_meta = root.find('.//main:futureMetadata', namespaces=NS)
                if future_meta is not None:
                    for bk in future_meta.findall('.//main:bk', namespaces=NS):
                        future_metadata_bk.append(bk)
            rv_list = []
            if 'xl/richData/rdrichvalue.xml' in namelist:
                root = ET.fromstring(z.read('xl/richData/rdrichvalue.xml'))
                for rv in root.findall('.//xlrd:rv', namespaces=NS):
                    v_tags = rv.findall('.//xlrd:v', namespaces=NS)
                    rv_list.append(int(v_tags[0].text) if v_tags else None)
            rel_list = []
            if 'xl/richData/richValueRel.xml' in namelist:
                root = ET.fromstring(z.read('xl/richData/richValueRel.xml'))
                for rel in root.findall('.//xlrr:rel', namespaces=NS):
                    rel_list.append(rel.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'])
            rv_targets = {}
            if 'xl/richData/_rels/richValueRel.xml.rels' in namelist:
                root = ET.fromstring(z.read('xl/richData/_rels/richValueRel.xml.rels'))
                for rel in root.findall('.//r:Relationship', namespaces={'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}):
                    rv_targets[rel.attrib['Id']] = rel.attrib['Target']
            images_to_insert = []
            for sheet_name, target in sheet_info.items():
                if target not in namelist: continue
                root = ET.fromstring(z.read(target))
                for c in root.findall('.//main:c', namespaces=NS):
                    if 'vm' in c.attrib:
                        vm_idx = int(c.attrib['vm']) - 1
                        if vm_idx < len(rc_list):
                            v_idx = rc_list[vm_idx]
                            if v_idx < len(future_metadata_bk):
                                rvb = None
                                for elem in bk.iter():
                                    if elem.tag.endswith('}rvb'):
                                        rvb = elem
                                        break
                                if rvb is not None:
                                    rv_idx = int(rvb.attrib.get('i', 0))
                                    if rv_idx < len(rv_list):
                                        local_img_id = rv_list[rv_idx]
                                        if local_img_id is not None and local_img_id < len(rel_list):
                                            r_id = rel_list[local_img_id]
                                            img_target = rv_targets.get(r_id)
                                            if img_target:
                                                img_path = img_target.replace('../', 'xl/')
                                                if img_path in namelist:
                                                    images_to_insert.append((sheet_name, c.attrib['r'], img_path))
        if not images_to_insert:
            return content_bytes
        wb = openpyxl.load_workbook(input_path)
        with zipfile.ZipFile(input_path, 'r') as z:
            for sheet_name, cell_ref, img_path in images_to_insert:
                ws = wb[sheet_name]
                ws[cell_ref].value = None
                img_temp_path = os.path.join(temp_dir, os.path.basename(img_path) + str(hash(cell_ref)) + ".png")
                with open(img_temp_path, 'wb') as f:
                    f.write(z.read(img_path))
                img = OpenpyxlImage(img_temp_path)
                img.width = 90
                img.height = 90
                ws.add_image(img, cell_ref)
                col_letter = ''.join([char for char in cell_ref if char.isalpha()])
                row_number = ''.join([char for char in cell_ref if char.isdigit()])
                ws.column_dimensions[col_letter].width = 12
                ws.row_dimensions[int(row_number)].height = 70
        wb.save(output_path)
        with open(output_path, 'rb') as f:
            return f.read()
    except Exception as e:
        print(f"Error repairing zoho sheet: {e}")
        return None
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

# ---------------------------------------------------------
# VERSION HISTORY ENDPOINTS
# ---------------------------------------------------------

@router.get("/{doc_id}/versions", response_model=list[VersionResponse])
async def get_sheet_versions(doc_id: str, limit: int = 50, offset: int = 0, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(SheetVersion)
        .where(SheetVersion.document_id == doc_id)
        .order_by(SheetVersion.created_at.desc())
        .limit(limit)
        .offset(offset)
    )
    return result.scalars().all()

@router.get("/{doc_id}/versions/{version_id}")
async def get_sheet_version_snapshot(doc_id: str, version_id: str, db: AsyncSession = Depends(get_db)):
    version = await db.get(SheetVersion, version_id)
    if not version or version.document_id != doc_id:
        raise HTTPException(404, "Version not found")
        
    from app.lib.document_storage import _s3_client, R2_BUCKET_NAME
    import gzip
    try:
        response = _s3_client.get_object(Bucket=R2_BUCKET_NAME, Key=version.sheet_snapshot_url)
        body_bytes = response["Body"].read()
        content = gzip.decompress(body_bytes).decode("utf-8")
        return {"content": json.loads(content)}
    except Exception as e:
        raise HTTPException(500, f"Failed to fetch snapshot: {str(e)}")

@router.post("/{doc_id}/versions")
async def create_named_version(
    doc_id: str, 
    body: VersionCreateRequest, 
    db: AsyncSession = Depends(get_db), 
    current_user: User = Depends(get_current_user)
):
    doc = await db.get(Document, doc_id)
    if not doc: raise HTTPException(404, "Not found")
    
    import main
    state = main.manager.doc_states.get(doc_id)
    if not state or not state.get("data"):
        raise HTTPException(400, "Document state is empty or not in memory")
        
    content_str = json.dumps(state["data"][0]) if isinstance(state["data"], list) and state["data"] else json.dumps(state["data"])
    
    import gzip
    from app.lib.document_storage import _s3_client, R2_BUCKET_NAME
    version_id = str(uuid.uuid4())
    sheet_snapshot_url = f"Drive/Sheet/{doc.owner_id}/versions/{doc_id}_{version_id}.json"
    body_bytes = gzip.compress(content_str.encode("utf-8"), compresslevel=1)
    
    import asyncio
    def upload_to_r2():
        _s3_client.put_object(
            Bucket=R2_BUCKET_NAME, Key=sheet_snapshot_url, Body=body_bytes, ContentType="application/json",
        )
    await asyncio.to_thread(upload_to_r2)
    
    IST_OFFSET = _dt.timezone(_dt.timedelta(hours=5, minutes=30))
    now = _dt.datetime.now(IST_OFFSET).replace(tzinfo=None)
    
    new_version = SheetVersion(
        id=version_id,
        document_id=doc_id,
        version_name=body.version_name,
        created_by_user_id=current_user.id,
        contributors=[current_user.id],
        is_named=True,
        sheet_snapshot_url=sheet_snapshot_url,
        created_at=now
    )
    db.add(new_version)
    await db.commit()
    return {"status": "success", "id": version_id}

@router.post("/{doc_id}/versions/{version_id}/restore")
async def restore_sheet_version(
    doc_id: str, 
    version_id: str, 
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    version = await db.get(SheetVersion, version_id)
    if not version or version.document_id != doc_id:
        raise HTTPException(404, "Version not found")
        
    doc = await db.get(Document, doc_id)
    from app.lib.document_storage import DocumentStorage, _s3_client, R2_BUCKET_NAME
    import gzip
    
    response = _s3_client.get_object(Bucket=R2_BUCKET_NAME, Key=version.sheet_snapshot_url)
    snapshot_bytes = response["Body"].read()
    snapshot_content = gzip.decompress(snapshot_bytes).decode("utf-8")
    
    new_version_id = str(uuid.uuid4())
    new_sheet_snapshot_url = f"Drive/Sheet/{doc.owner_id}/versions/{doc_id}_{new_version_id}.json"
    
    _s3_client.put_object(
        Bucket=R2_BUCKET_NAME, Key=new_sheet_snapshot_url, Body=snapshot_bytes, ContentType="application/json",
    )
    
    IST_OFFSET = _dt.timezone(_dt.timedelta(hours=5, minutes=30))
    now = _dt.datetime.now(IST_OFFSET).replace(tzinfo=None)
    
    new_version = SheetVersion(
        id=new_version_id,
        document_id=doc_id,
        version_name=f"Restored from {version.created_at.strftime('%b %d, %H:%M')}",
        created_by_user_id=current_user.id,
        contributors=[current_user.id],
        is_named=True,
        sheet_snapshot_url=new_sheet_snapshot_url,
        base_version_id=version.id,
        created_at=now
    )
    db.add(new_version)
    
    storage_res = DocumentStorage.save(doc.owner_id, doc.id, snapshot_content, doc_type=doc.doc_type)
    doc.file_path = storage_res["relative_path"]
    doc.file_size = storage_res["size"]
    doc.content_version = (doc.content_version or 1) + 1
    
    await db.commit()
    
    import main
    if doc_id in main.manager.doc_states:
        main.manager.doc_states.pop(doc_id, None)
        
    return {"status": "success", "id": new_version_id}

@router.post("/import")
async def import_document(
    file: UploadFile = File(...),
    replace_doc_id: Optional[str] = Form(None),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    filename = file.filename
    content_bytes = await file.read()
    
    doc_type = "doc"
    filename_lower = filename.lower()
    if filename_lower.endswith(".xlsx") or filename_lower.endswith(".xls") or filename_lower.endswith(".csv") or filename_lower.endswith(".tsv"):
        doc_type = "sheet"
    elif filename_lower.endswith(".pptx"):
        doc_type = "slide"
        
    title = filename.rsplit(".", 1)[0]
    content_json = "{}"
    
    if doc_type == "sheet" and (filename_lower.endswith(".xlsx") or filename_lower.endswith(".xls")):
        try:
            if filename_lower.endswith(".xlsx"):
                content_bytes = repair_zoho_rich_values(content_bytes)
                
            from openpyxl import load_workbook
            from openpyxl.utils import column_index_from_string
            wb = load_workbook(io.BytesIO(content_bytes), data_only=True)
            wb_formula = load_workbook(io.BytesIO(content_bytes), data_only=False)
            sheets = []
            
            for ws_idx, ws in enumerate(wb.worksheets):
                ws_f = wb_formula.worksheets[ws_idx]
                ROWS = max(1000, ws.max_row + 5)
                COLS = max(26, ws.max_column + 5)
                
                cells = [["" for _ in range(COLS)] for _ in range(ROWS)]
                formats = {}
                validations = {}
                colWidths = {}
                rowHeights = {}
                
                for row in ws.iter_rows():
                    for cell in row:
                        val = cell.value
                        if val is None or (isinstance(val, str) and val.startswith('#')):
                            f_cell = ws_f.cell(row=cell.row, column=cell.column)
                            if isinstance(f_cell.value, str) and f_cell.value.upper().startswith('=IMAGE('):
                                val = f_cell.value

                        if val is not None:
                            r = cell.row - 1
                            c = cell.column - 1
                            if r < ROWS and c < COLS:
                                if isinstance(val, (_dt.datetime, _dt.date)):
                                    cells[r][c] = val.strftime("%d/%m/%Y")
                                else:
                                    cells[r][c] = str(val)
                            
                            fmt = {}
                            if cell.font:
                                if cell.font.bold: fmt["bold"] = True
                                if cell.font.italic: fmt["italic"] = True
                                if cell.font.strike: fmt["strikethrough"] = True
                                try:
                                    rgb = cell.font.color.rgb
                                    if isinstance(rgb, str) and len(rgb) in [6, 8]:
                                        if len(rgb) == 8: fmt["color"] = f"#{rgb[2:]}"
                                        elif len(rgb) == 6: fmt["color"] = f"#{rgb}"
                                except:
                                    pass
                            if cell.alignment:
                                if cell.alignment.wrapText: fmt["wrap"] = True
                            if cell.fill and cell.fill.fill_type and cell.fill.fill_type != 'none':
                                fgColor = cell.fill.fgColor
                                if fgColor:
                                    rgb = None
                                    if fgColor.type == 'rgb':
                                        rgb = fgColor.rgb
                                    elif hasattr(fgColor, 'value') and isinstance(fgColor.value, str) and len(fgColor.value) in [6, 8]:
                                        rgb = fgColor.value
                                    if rgb and isinstance(rgb, str):
                                        if len(rgb) == 8: fmt["bg"] = f"#{rgb[2:]}"
                                        elif len(rgb) == 6: fmt["bg"] = f"#{rgb}"
                            if fmt:
                                formats[f"{r},{c}"] = fmt
                                
                if ws.merged_cells:
                    for merged_range in ws.merged_cells.ranges:
                        min_col, min_row, max_col, max_row = merged_range.bounds
                        min_r, min_c = min_row - 1, min_col - 1
                        max_r, max_c = max_row - 1, max_col - 1
                        
                        tl_ref = f"{min_r},{min_c}"
                        if tl_ref not in formats:
                            formats[tl_ref] = {}
                        formats[tl_ref]["_mergeSpan"] = {
                            "rows": max_r - min_r + 1,
                            "cols": max_c - min_c + 1
                        }
                        
                        for r in range(min_r, max_r + 1):
                            for c in range(min_c, max_c + 1):
                                if r == min_r and c == min_c:
                                    continue
                                slave_ref = f"{r},{c}"
                                if slave_ref not in formats:
                                    formats[slave_ref] = {}
                                formats[slave_ref]["_mergedInto"] = tl_ref

                if hasattr(ws, '_images') and ws._images:
                    import base64
                    for img in ws._images:
                        try:
                            anchor = img.anchor
                            if hasattr(anchor, '_from'):
                                row_idx = anchor._from.row
                                col_idx = anchor._from.col
                            elif hasattr(anchor, 'row') and hasattr(anchor, 'col'):
                                row_idx = anchor.row - 1
                                col_idx = anchor.col - 1
                            else:
                                continue
                            
                            # Get image format (default to png)
                            fmt = getattr(img, 'format', 'png')
                            if not fmt:
                                fmt = 'png'
                            
                            # Get bytes
                            img_bytes = img._data()
                            b64_data = base64.b64encode(img_bytes).decode('utf-8')
                            mime = f"image/{fmt.lower()}"
                            if row_idx < ROWS and col_idx < COLS:
                                cells[row_idx][col_idx] = f"data:{mime};base64,{b64_data}"
                        except Exception as img_err:
                            try:
                                with open("backend/debug_validation.log", "a") as f_log:
                                    f_log.write(f"  Error parsing image: {img_err}\n")
                            except:
                                pass

                if ws.data_validations:
                    for dv in ws.data_validations.dataValidation:
                        if dv.type == "list":
                            options = resolve_validation_options(dv.formula1, wb, ws)
                            if options:
                                sqref = dv.sqref
                                ranges_list = []
                                if isinstance(sqref, str):
                                    for ref_part in sqref.split(' '):
                                        try:
                                            from openpyxl.utils.cell import range_boundaries
                                            min_col, min_row, max_col, max_row = range_boundaries(ref_part)
                                            ranges_list.append((min_row, max_row, min_col, max_col))
                                        except:
                                            pass
                                elif hasattr(sqref, 'ranges'):
                                    for r_obj in sqref.ranges:
                                        ranges_list.append((r_obj.min_row, r_obj.max_row, r_obj.min_col, r_obj.max_col))
                                else:
                                    # Fallback
                                    try:
                                        sqref_str = str(sqref)
                                        for ref_part in sqref_str.split(' '):
                                            from openpyxl.utils.cell import range_boundaries
                                            min_col, min_row, max_col, max_row = range_boundaries(ref_part)
                                            ranges_list.append((min_row, max_row, min_col, max_col))
                                    except:
                                        pass
                                        
                                for min_row, max_row, min_col, max_col in ranges_list:
                                    for r in range(min_row - 1, max_row):
                                        for c in range(min_col - 1, max_col):
                                            if r < ROWS and c < COLS:
                                                validations[f"{r},{c}"] = {"type": "list", "options": options}
                
                for col_letter, col_dim in ws.column_dimensions.items():
                    try:
                        c = column_index_from_string(col_letter) - 1
                        if col_dim.width:
                            colWidths[str(c)] = int(col_dim.width * 7)
                    except:
                        pass
                        
                for row_num, row_dim in ws.row_dimensions.items():
                    r = row_num - 1
                    if row_dim.height:
                        rowHeights[str(r)] = int(row_dim.height * 1.3)
                            
                sparse_cells: dict = {}
                for ri, row_data in enumerate(cells):
                    for ci, v in enumerate(row_data):
                        if v:
                            if ri not in sparse_cells:
                                sparse_cells[ri] = {}
                            sparse_cells[ri][ci] = v

                sheets.append({
                    "name": ws.title,
                    "cells": sparse_cells,
                    "formats": formats,
                    "validations": validations,
                    "colWidths": colWidths,
                    "rowHeights": rowHeights
                })
                
            content_json = json.dumps({
                "_importedSheets": sheets
            })
        except Exception as e:
            import traceback
            print(f"[IMPORT XLSX ERROR] {e}")
            print(traceback.format_exc())
            content_json = _default_content(doc_type, title)
    elif doc_type == "sheet" and (filename_lower.endswith(".csv") or filename_lower.endswith(".tsv")):
        try:
            import csv as _csv
            delimiter = "\t" if filename_lower.endswith(".tsv") else ","
            text = content_bytes.decode("utf-8-sig", errors="replace")
            reader = _csv.reader(text.splitlines(), delimiter=delimiter)
            rows_data = list(reader)
            sparse_cells = {}
            for ri, row in enumerate(rows_data):
                for ci, val in enumerate(row):
                    if val.strip():
                        if ri not in sparse_cells:
                            sparse_cells[ri] = {}
                        sparse_cells[ri][ci] = val
            content_json = json.dumps({
                "_importedSheets": [{
                    "name": title,
                    "cells": sparse_cells,
                    "formats": {},
                    "validations": {},
                    "colWidths": {},
                    "rowHeights": {}
                }]
            })
        except Exception as e:
            import traceback
            print(f"[IMPORT CSV ERROR] {e}")
            print(traceback.format_exc())
            content_json = _default_content(doc_type, title)
    elif doc_type == "doc":
        parsed = False
        parse_error = None
        try:
            import docx
            import base64
            doc_file = docx.Document(io.BytesIO(content_bytes))
            html = ""

            def _run_to_html(run):
                """Convert a single docx run to HTML, handling bold/italic/underline."""
                text = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                if not text:
                    return ""
                if run.bold:      text = f"<b>{text}</b>"
                if run.italic:    text = f"<i>{text}</i>"
                if run.underline: text = f"<u>{text}</u>"
                if run.font.strike: text = f"<s>{text}</s>"
                return text

            def _alignment_style(p):
                try:
                    if p.alignment == 1: return ' style="text-align:center;"'
                    if p.alignment == 2: return ' style="text-align:right;"'
                    if p.alignment == 3: return ' style="text-align:justify;"'
                except:
                    pass
                return ""

            def _para_tag(p):
                """Return the appropriate HTML tag for a paragraph's style."""
                try:
                    style = (p.style.name or "").lower()
                    if "heading 1" in style: return "h1"
                    if "heading 2" in style: return "h2"
                    if "heading 3" in style: return "h3"
                    if "heading 4" in style: return "h4"
                except:
                    pass
                return "p"

            # Process body elements (paragraphs + tables) in order
            from docx.oxml.ns import qn
            for element in doc_file.element.body:
                tag = element.tag.split("}")[-1]  # strip namespace

                if tag == "p":
                    # Wrap as a paragraph object
                    from docx.text.paragraph import Paragraph
                    p = Paragraph(element, doc_file)
                    p_html = ""

                    # Images
                    for run in p.runs:
                        drawings = run._element.xpath('.//w:drawing')
                        for drawing in drawings:
                            for blip in drawing.xpath('.//a:blip'):
                                embed_id = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                                if embed_id and embed_id in doc_file.part.related_parts:
                                    part = doc_file.part.related_parts[embed_id]
                                    b64 = base64.b64encode(part.blob).decode("utf-8")
                                    if p_html.strip():
                                        html += f"<{_para_tag(p)}{_alignment_style(p)}>{p_html}</{_para_tag(p)}>"
                                        p_html = ""
                                    html += f'<img src="data:{part.content_type};base64,{b64}" style="max-width:100%;height:auto;display:block;margin:12px auto;"/>'
                        p_html += _run_to_html(run)

                    if p_html.strip():
                        tag_name = _para_tag(p)
                        html += f"<{tag_name}{_alignment_style(p)}>{p_html}</{tag_name}>"
                    elif not p_html:
                        html += "<br>"

                elif tag == "tbl":
                    from docx.table import Table
                    tbl = Table(element, doc_file)
                    html += '<table border="1" style="border-collapse:collapse;width:100%;margin:8px 0;"><tbody>'
                    for row in tbl.rows:
                        html += "<tr>"
                        for cell in row.cells:
                            cell_text = cell.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            html += f'<td style="padding:4px;">{cell_text}</td>'
                        html += "</tr>"
                    html += "</tbody></table>"

            if not html:
                html = "<br>"
            content_json = json.dumps({"html": html})
            parsed = True
            print(f"[IMPORT] docx parsed OK: {len(html)} chars of HTML for doc '{title}'")
        except Exception as e:
            import traceback
            parse_error = traceback.format_exc()
            print(f"[IMPORT] docx parse FAILED for '{title}': {e}")
            print(parse_error)

        if not parsed:
            if filename_lower.endswith(".html"):
                try:
                    html_content = content_bytes.decode('utf-8', errors='replace')
                    content_json = json.dumps({"html": html_content})
                    parsed = True
                except:
                    pass
            elif filename_lower.endswith(".txt"):
                # Only decode as text if the file is actually a text file
                try:
                    text = content_bytes.decode('utf-8', errors='replace')
                    html_content = "".join([f"<p>{line}</p>" for line in text.splitlines()])
                    content_json = json.dumps({"html": html_content or "<br>"})
                    parsed = True
                except:
                    pass
            # NOTE: do NOT fall back to decoding .docx/.pdf binary bytes as UTF-8 text
            # that produces garbage. Just return a friendly error message.
            if not parsed:
                error_msg = f"Could not parse {filename}."
                if parse_error:
                    # Include first line of error for debug visibility in editor
                    first_line = parse_error.strip().splitlines()[-1]
                    error_msg += f" Server error: {first_line}"
                content_json = json.dumps({"html": f"<p style='color:red;'>{error_msg} Please try a .txt or .html file.</p>"})
                parsed = True

        if not parsed:
            content_json = _default_content(doc_type, title)
            
    elif doc_type == "slide":
        try:
            from pptx import Presentation
            import zipfile
            prs = Presentation(io.BytesIO(content_bytes))
            pages = {}
            page_order = []
            for slide in prs.slides:
                pid = str(uuid.uuid4())
                page_order.append(pid)
                text_content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += f"<p>{shape.text}</p>"
                pages[pid] = {"title": f"Slide {len(page_order)}", "body": text_content}
            if not pages:
                content_json = _default_content(doc_type, title)
            else:
                content_json = json.dumps({
                    "id": page_order[0] if page_order else str(uuid.uuid4()),
                    "title": title,
                    "pages": pages,
                    "pageOrder": page_order
                })
        except Exception as e:
            content_json = _default_content(doc_type, title)
    else:
        content_json = _default_content(doc_type, title)
        
    if replace_doc_id:
        doc = await db.get(Document, replace_doc_id)
        if not doc or doc.owner_id != current_user.id:
            raise HTTPException(404, "Document not found or access denied")
        doc.title = title
    else:
        doc = Document(
            title=title,
            doc_type=doc_type,
            owner_id=current_user.id,
        )
        db.add(doc)
    await db.commit()
    await db.refresh(doc)

    # Save to R2 storage in background — don't let storage failure block the response
    try:
        storage_res = await asyncio.to_thread(
            DocumentStorage.save, doc.owner_id, doc.id, content_json, doc_type=doc.doc_type
        )
        doc.file_path = storage_res["relative_path"]
        doc.file_size = storage_res["size"]
        doc.content_version = 1
        await db.commit()
        print(f"[IMPORT] Saved doc {doc.id} to R2 ({storage_res['size']} bytes)")
    except Exception as e:
        import traceback
        print(f"[IMPORT] R2 save failed for doc {doc.id}: {e}")
        print(traceback.format_exc())
        # Still return success — frontend uses response content directly

    # Return the parsed content directly so frontend can apply it without a reload
    return {"id": doc.id, "title": doc.title, "doc_type": doc.doc_type, "content": content_json}


@router.get("")
async def list_documents(
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    # Select only the columns we need — NOT content (can be 50 MB+)
    result = await db.execute(
        select(
            Document.id,
            Document.title,
            Document.doc_type,
            Document.updated_at,
            Document.is_trashed,
            Document.owner_id,
            User.name.label('owner_name'),
        )
        .join(User, Document.owner_id == User.id)
        .outerjoin(DocumentShare, Document.id == DocumentShare.document_id)
        .where(
            or_(
                Document.owner_id == current_user.id,
                DocumentShare.user_id == current_user.id,
            ),
        )
        .order_by(Document.updated_at.desc())
        .distinct()
    )
    docs = result.all()
    return [
        {
            "id": row.id,
            "title": row.title,
            "doc_type": row.doc_type,
            "updated_at": row.updated_at,
            "is_trashed": row.is_trashed,
            "owner_id": row.owner_id,
            "owner_name": row.owner_name,
        }
        for row in docs
    ]


@router.get("/{doc_id}")
async def get_document(
    doc_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[User] = Depends(get_optional_current_user),
):
    doc = await db.get(Document, doc_id)
    if not doc:
        raise HTTPException(404, "Not found")

    if not doc.is_public:
        if not current_user:
            raise HTTPException(401, "Not authenticated")
        if doc.owner_id != current_user.id:
            existing = await db.execute(
                select(DocumentShare).where(
                    DocumentShare.document_id == doc_id,
                    DocumentShare.user_id == current_user.id,
                )
            )
            if not existing.scalar_one_or_none():
                raise HTTPException(403, "Access denied")

    content = "{}"
    if doc.file_path:
        try:
            content = await asyncio.to_thread(
                DocumentStorage.load,
                doc.owner_id, doc.id,
                doc_type=doc.doc_type,
                file_path=doc.file_path,
            )
        except Exception as e:
            import traceback
            print(f"Failed to load doc {doc.id} from storage: {e}")
            print(traceback.format_exc())
            pass

    return {"id": doc.id, "title": doc.title, "doc_type": doc.doc_type, "content": content}


@router.get("/{doc_id}/audit-events")
async def get_audit_events(
    doc_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    doc = await db.get(Document, doc_id)
    if not doc:
        raise HTTPException(404, "Not found")

    result = await db.execute(
        select(AuditEvent)
        .where(AuditEvent.document_id == doc_id)
        .order_by(AuditEvent.created_at.desc())
        .limit(1000)
    )
    events = result.scalars().all()
    return events


@router.post("/{doc_id}/audit-events")
async def save_audit_events(
    doc_id: str,
    events: list[dict],
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    doc = await db.get(Document, doc_id)
    if not doc:
        raise HTTPException(404, "Not found")

    for ev in events:
        ae = AuditEvent(
            document_id=doc_id,
            user_name=current_user.name or 'unknown',
            sheet_id=ev.get("sheet_id", ""),
            sheet_name=ev.get("sheet_name", ""),
            action_type=ev.get("action_type", ""),
            target_range=ev.get("target_range", ""),
            metadata_json=ev.get("metadata"),
        )
        db.add(ae)
        
    # VERSION HISTORY AUTO-CHECKPOINT LOGIC
    IST_OFFSET = _dt.timezone(_dt.timedelta(hours=5, minutes=30))
    now = _dt.datetime.now(IST_OFFSET).replace(tzinfo=None)
    
    last_version_result = await db.execute(
        select(SheetVersion)
        .where(SheetVersion.document_id == doc_id, SheetVersion.is_named == False)
        .order_by(SheetVersion.created_at.desc())
        .limit(1)
    )
    last_version = last_version_result.scalar_one_or_none()
    
    needs_checkpoint = False
    if not last_version:
        needs_checkpoint = True
    elif (now - last_version.created_at).total_seconds() > 600:
        needs_checkpoint = True
        
    if needs_checkpoint:
        import main
        state = main.manager.doc_states.get(doc_id)
        if state and state.get("data"):
            content_str = json.dumps(state["data"][0]) if isinstance(state["data"], list) and state["data"] else json.dumps(state["data"])
            
            import gzip
            from app.lib.document_storage import _s3_client, R2_BUCKET_NAME
            version_id = str(uuid.uuid4())
            sheet_snapshot_url = f"Drive/Sheet/{doc.owner_id}/versions/{doc_id}_{version_id}.json"
            body_bytes = gzip.compress(content_str.encode("utf-8"), compresslevel=1)
            
            import asyncio
            def upload_to_r2():
                _s3_client.put_object(
                    Bucket=R2_BUCKET_NAME,
                    Key=sheet_snapshot_url,
                    Body=body_bytes,
                    ContentType="application/json",
                )
            await asyncio.to_thread(upload_to_r2)
            
            contributors = [current_user.id]
            
            new_version = SheetVersion(
                id=version_id,
                document_id=doc_id,
                created_by_user_id=current_user.id,
                contributors=contributors,
                is_named=False,
                sheet_snapshot_url=sheet_snapshot_url,
                created_at=now
            )
            db.add(new_version)

    await db.commit()
    return {"status": "success"}


@router.put("/{doc_id}")
async def update_document(
    doc_id: str,
    body: DocumentUpdate,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    doc = await db.get(Document, doc_id)
    if not doc:
        raise HTTPException(404, "Not found")

    # Allow owner or users with edit permission
    if doc.owner_id != current_user.id:
        share = await db.execute(
            select(DocumentShare).where(
                DocumentShare.document_id == doc_id,
                DocumentShare.user_id == current_user.id,
                DocumentShare.permission.in_(["edit", "Edit", "EDIT"]),
            )
        )
        if not share.scalar_one_or_none():
            raise HTTPException(403, "Forbidden")

    if body.title is not None:
        doc.title = body.title
    if body.content is not None:
        # SAFETY: Refuse to overwrite substantial existing content with empty data.
        # This prevents the frontend race condition where empty content is saved
        # before the real document data has finished loading from R2.
        skip_save = False
        if doc.doc_type == "sheet" and doc.file_size and doc.file_size > 500:
            try:
                parsed_new = json.loads(body.content)
                has_cell_data = False
                sheets = parsed_new.get("_importedSheets", [])
                if not sheets and parsed_new.get("cells"):
                    has_cell_data = any(parsed_new["cells"].values()) if isinstance(parsed_new["cells"], dict) else bool(parsed_new["cells"])
                for s in sheets:
                    sc = s.get("cells", {})
                    if isinstance(sc, dict) and sc:
                        has_cell_data = True
                        break
                    elif isinstance(sc, list) and any(any(cell for cell in row) for row in sc if row):
                        has_cell_data = True
                        break
                if not has_cell_data:
                    skip_save = True
                    print(f"[SAVE GUARD] Blocked saving empty sheet content for doc {doc.id} (existing size: {doc.file_size} bytes)")
            except Exception:
                pass  # If we can't parse, allow the save

        # SAFETY: For doc type, block overwriting real content with empty/trivial HTML.
        # Threshold: existing file > 500 bytes, new html content is empty or just <br>/<div><br></div>.
        if doc.doc_type == "doc" and doc.file_size and doc.file_size > 500:
            try:
                parsed_new = json.loads(body.content)
                new_html = parsed_new.get("html", "")
                stripped = new_html.strip().replace("<br>", "").replace("<div>", "").replace("</div>", "").strip()
                if not stripped:
                    skip_save = True
                    print(f"[SAVE GUARD] Blocked saving blank doc content for doc {doc.id} (existing size: {doc.file_size} bytes)")
            except Exception:
                pass  # If we can't parse, allow the save

        if not skip_save:
            try:
                storage_res = await asyncio.to_thread(
                    DocumentStorage.save, doc.owner_id, doc.id, body.content, doc_type=doc.doc_type
                )
                doc.file_path = storage_res["relative_path"]
                doc.file_size = storage_res["size"]
                doc.content_version = (doc.content_version or 1) + 1
                
                # Checkpoint logic for manual/explicit saves (debounce 30 seconds)
                if doc.doc_type == "sheet":
                    IST_OFFSET = _dt.timezone(_dt.timedelta(hours=5, minutes=30))
                    now = _dt.datetime.now(IST_OFFSET).replace(tzinfo=None)
                    
                    last_version_result = await db.execute(
                        select(SheetVersion)
                        .where(SheetVersion.document_id == doc_id, SheetVersion.is_named == False)
                        .order_by(SheetVersion.created_at.desc())
                        .limit(1)
                    )
                    last_version = last_version_result.scalar_one_or_none()
                    
                    if not last_version or (now - last_version.created_at).total_seconds() > 30:
                        version_id = str(uuid.uuid4())
                        sheet_snapshot_url = f"Drive/Sheet/{doc.owner_id}/versions/{doc_id}_{version_id}.json"
                        
                        import gzip
                        from app.lib.document_storage import _s3_client, R2_BUCKET_NAME
                        body_bytes = gzip.compress(body.content.encode("utf-8"), compresslevel=1)
                        
                        def upload_to_r2():
                            _s3_client.put_object(
                                Bucket=R2_BUCKET_NAME,
                                Key=sheet_snapshot_url,
                                Body=body_bytes,
                                ContentType="application/json",
                            )
                        await asyncio.to_thread(upload_to_r2)
                        
                        new_version = SheetVersion(
                            id=version_id,
                            document_id=doc_id,
                            created_by_user_id=current_user.id,
                            contributors=[current_user.id],
                            is_named=False,
                            sheet_snapshot_url=sheet_snapshot_url,
                            created_at=now
                        )
                        db.add(new_version)
                        
            except Exception as e:
                import traceback
                print(f"Failed to save document {doc.id} to storage: {e}")
                print(traceback.format_exc())
            
    if body.is_trashed is not None:
        doc.is_trashed = body.is_trashed
        
    await db.commit()
    return {"id": doc.id, "title": doc.title}


@router.delete("/{doc_id}")
async def delete_document(
    doc_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    doc = await db.get(Document, doc_id)
    if not doc:
        raise HTTPException(404, "Not found")
    if doc.owner_id != current_user.id:
        raise HTTPException(403, "Only the owner can delete this document")
    
    if doc.is_trashed:
        try:
            await asyncio.to_thread(
                DocumentStorage.delete,
                doc.owner_id, doc.id,
                doc_type=doc.doc_type,
                file_path=doc.file_path,
            )
        except Exception as e:
            print(f"Failed to delete document {doc.id} from storage: {e}")
        await db.delete(doc)
    else:
        doc.is_trashed = 1
        
    await db.commit()
    return {"deleted": True}

@router.post("/{doc_id}/share")
async def share_document(
    doc_id: str,
    body: DocumentShareCreate,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):

    doc = await db.get(Document, doc_id)
    if not doc or doc.owner_id != current_user.id:
        raise HTTPException(404, "Not found or not owner")

    # Try to find user by email in the vmail users table
    result = await db.execute(select(User).where(User.email == body.email))
    target_user = result.scalar_one_or_none()

    if target_user:
        # --- Internal user: share via user_id as before ---
        existing = await db.execute(select(DocumentShare).where(
            DocumentShare.document_id == doc_id,
            DocumentShare.user_id == target_user.id
        ))
        share = existing.scalar_one_or_none()
        if share:
            share.permission = body.permission
        else:
            share = DocumentShare(
                document_id=doc_id,
                user_id=target_user.id,
                permission=body.permission,
            )
            db.add(share)
        await db.commit()
        # Send internal notification email
        _send_share_email(
            to_email=body.email,
            to_name=target_user.name,
            from_name=current_user.name,
            from_email=current_user.email,
            doc_title=doc.title,
            doc_type=doc.doc_type,
            doc_id=doc_id,
            permission=body.permission,
            is_external=False,
        )
        return {"shared": True, "permission": body.permission, "user_id": target_user.id, "external": False}
    else:
        # --- External user: share via external_email, no user_id ---
        existing = await db.execute(select(DocumentShare).where(
            DocumentShare.document_id == doc_id,
            DocumentShare.external_email == body.email,
        ))
        share = existing.scalar_one_or_none()
        if share:
            share.permission = body.permission
        else:
            # Also make the doc public so the link is accessible without login
            doc.is_public = True
            share = DocumentShare(
                document_id=doc_id,
                user_id=None,
                external_email=body.email,
                permission=body.permission,
            )
            db.add(share)
        await db.commit()
        # Send external notification email
        _send_share_email(
            to_email=body.email,
            to_name=body.email.split("@")[0],
            from_name=current_user.name,
            from_email=current_user.email,
            doc_title=doc.title,
            doc_type=doc.doc_type,
            doc_id=doc_id,
            permission=body.permission,
            is_external=True,
        )
        return {"shared": True, "permission": body.permission, "user_id": None, "external": True}


def _send_share_email(
    to_email: str,
    to_name: str,
    from_name: str,
    from_email: str,
    doc_title: str,
    doc_type: str,
    doc_id: str,
    permission: str,
    is_external: bool,
):
    """Send a share notification email via SMTP. Silently fails if SMTP not configured."""
    smtp_host     = os.getenv("SMTP_HOST", "")
    smtp_port     = int(os.getenv("SMTP_PORT", "587"))
    smtp_user     = os.getenv("SMTP_USER", "")
    smtp_password = os.getenv("SMTP_PASSWORD", "")
    smtp_from     = os.getenv("SMTP_FROM", smtp_user)
    frontend_url  = os.getenv("SHEETS_URL", os.getenv("FRONTEND_URL", "https://sheets.vsnaptechnology.com"))

    if not smtp_host or not smtp_user or not smtp_password:
        print(f"[SHARE EMAIL] SMTP not configured — skipping email to {to_email}")
        return

    # Build the doc URL based on doc_type
    if doc_type == "sheet":
        base_url = os.getenv("SHEETS_URL", "https://sheets.vsnaptechnology.com")
        doc_url = f"{base_url}/sheet/{doc_id}"
    elif doc_type == "doc":
        base_url = os.getenv("DOCS_URL", "https://docs.vsnaptechnology.com")
        doc_url = f"{base_url}/doc/{doc_id}"
    else:
        base_url = os.getenv("SHOW_URL", "https://show.vsnaptechnology.com")
        doc_url = f"{base_url}/slide/{doc_id}"

    subject = f"{from_name} has shared a file with you"

    type_label = {"sheet": "Spreadsheet", "doc": "Document", "slide": "Presentation"}.get(doc_type, "File")

    html_body = f"""
<!DOCTYPE html>
<html>
<head><meta charset="UTF-8"></head>
<body style="font-family: Arial, sans-serif; background:#f5f5f5; margin:0; padding:20px;">
  <table width="100%" cellpadding="0" cellspacing="0">
    <tr><td align="center">
      <table width="560" style="background:#ffffff; border-radius:8px; padding:32px; border:1px solid #e0e0e0;">
        <tr>
          <td style="padding-bottom:24px; border-bottom:1px solid #eeeeee;">
            <span style="font-size:22px; font-weight:700; color:#1a73e8;">VSNap Office Suite</span>
          </td>
        </tr>
        <tr>
          <td style="padding-top:24px; padding-bottom:16px; font-size:15px; color:#333;">
            Hello <strong>{to_name}</strong>,
          </td>
        </tr>
        <tr>
          <td style="padding-bottom:24px; font-size:15px; color:#333; line-height:1.6;">
            <strong>{from_name}</strong>
            (<a href="mailto:{from_email}" style="color:#1a73e8; text-decoration:none;">{from_email}</a>)
            {'— External user ' if is_external else ''}has shared the {type_label}
            <strong>"{doc_title}"</strong> with you with
            <strong>{permission.capitalize()}</strong> permission.
          </td>
        </tr>
        <tr>
          <td align="center" style="padding-bottom:32px;">
            <a href="{doc_url}" style="background:#1a73e8; color:#fff; text-decoration:none; padding:12px 32px; border-radius:4px; font-size:15px; font-weight:500; display:inline-block;">
              View {type_label}
            </a>
          </td>
        </tr>
        <tr>
          <td style="font-size:12px; color:#999; border-top:1px solid #eeeeee; padding-top:16px;">
            This email was sent by VSNap Office Suite on behalf of {from_name}.
          </td>
        </tr>
      </table>
    </td></tr>
  </table>
</body>
</html>
"""

    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = subject
        msg["From"]    = f"{from_name} (via VSNap Office Suite) <{smtp_from}>"
        msg["To"]      = to_email
        msg.attach(MIMEText(html_body, "html"))

        with smtplib.SMTP(smtp_host, smtp_port) as server:
            server.ehlo()
            server.starttls()
            server.login(smtp_user, smtp_password)
            server.sendmail(smtp_from, [to_email], msg.as_string())
        print(f"[SHARE EMAIL] Sent to {to_email} successfully")
    except Exception as e:
        print(f"[SHARE EMAIL] Failed to send to {to_email}: {e}")