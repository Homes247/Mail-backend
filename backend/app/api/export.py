import io, json
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.ext.asyncio import AsyncSession
from pydantic import BaseModel
from app.database import get_db
from app.models.document import Document

router = APIRouter()

class ExportRequest(BaseModel):
    doc_id: str
    format: str  # csv | docx | pptx

@router.post("/export")
async def export_document(body: ExportRequest, db: AsyncSession = Depends(get_db)):
    doc = await db.get(Document, body.doc_id)
    if not doc:
        raise HTTPException(404, "Not found")

    content = json.loads(doc.content or "{}")

    if body.format in ["csv", "tsv"]:
        return _export_csv(doc.title, content, delimiter='\t' if body.format == 'tsv' else ',')
    elif body.format in ["xlsx", "xlsb", "ods"]:
        return _export_xlsx(doc.title, content)
    elif body.format == "html":
        return _export_html_zip(doc.title, content)
    elif body.format == "pdf":
        return _export_pdf(doc.title, content)
    elif body.format == "docx":
        return _export_docx(doc.title, content)
    elif body.format == "pptx":
        return _export_pptx(doc.title, content)
    else:
        raise HTTPException(400, "Unsupported format")


def _export_csv(title: str, content: dict, delimiter: str = ',') -> StreamingResponse:
    import csv, io
    output = io.StringIO()
    writer = csv.writer(output, delimiter=delimiter)

    cells = content.get("cells", {})
    if cells:
        max_row = max((int(r) for r in cells.keys()), default=-1)
        max_col = max((int(c) for r in cells.values() for c in r.keys()), default=-1)
        
        for r in range(max_row + 1):
            row_data = []
            for c in range(max_col + 1):
                val = cells.get(str(r), {}).get(str(c), "")
                row_data.append(val)
            writer.writerow(row_data)

    buf = io.BytesIO(output.getvalue().encode('utf-8'))
    ext = "tsv" if delimiter == '\t' else "csv"
    return StreamingResponse(
        buf,
        media_type="text/csv",
        headers={"Content-Disposition": f'attachment; filename="{title}.{ext}"'}
    )

def _export_xlsx(title: str, content: dict) -> StreamingResponse:
    import io
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    cells = content.get("cells", {})
    if cells:
        for r_str, cols in cells.items():
            r = int(r_str) + 1  # 1-indexed for excel
            for c_str, val in cols.items():
                c = int(c_str) + 1
                try:
                    ws.cell(row=r, column=c, value=val)
                except Exception:
                    ws.cell(row=r, column=c, value=str(val))

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{title}.xlsx"'}
    )

def _export_html_zip(title: str, content: dict) -> StreamingResponse:
    import io, zipfile
    html_parts = [f"<html><head><title>{title}</title><style>table {{ border-collapse: collapse; }} td {{ border: 1px solid #ddd; padding: 4px; }}</style></head><body>"]
    html_parts.append(f"<h2>{title}</h2><table>")
    
    cells = content.get("cells", {})
    if cells:
        max_row = max((int(r) for r in cells.keys()), default=-1)
        max_col = max((int(c) for r in cells.values() for c in r.keys()), default=-1)
        for r in range(max_row + 1):
            html_parts.append("<tr>")
            for c in range(max_col + 1):
                val = cells.get(str(r), {}).get(str(c), "")
                html_parts.append(f"<td>{val}</td>")
            html_parts.append("</tr>")
            
    html_parts.append("</table></body></html>")
    html_str = "".join(html_parts)
    
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{title}.html", html_str.encode("utf-8"))
    
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{title}_html.zip"'}
    )



def _export_docx(title: str, content: dict) -> StreamingResponse:
    from docx import Document as DocxDocument
    import re

    doc = DocxDocument()
    doc.add_heading(title, 0)

    # New format: { html: "<p>...</p>" }
    html = content.get("html", "")
    if html:
        # Strip tags and split into paragraphs for basic export
        clean = re.sub(r'<br\s*/?>', '\n', html)
        clean = re.sub(r'<[^>]+>', '', clean)
        for line in clean.split('\n'):
            stripped = line.strip()
            if stripped:
                doc.add_paragraph(stripped)
    else:
        doc.add_paragraph("(empty document)")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{title}.docx"'}
    )


def _export_pptx(title: str, content: dict) -> StreamingResponse:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title + Content

    pages = content.get("pages", {})
    page_order = content.get("pageOrder", list(pages.keys()))

    if not page_order:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    else:
        for page_id in page_order:
            page = pages.get(page_id, {})
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = page.get("title", "")
            # Write body text into content placeholder
            body_text = page.get("body", "")
            if body_text and len(slide.placeholders) > 1:
                slide.placeholders[1].text = body_text

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{title}.pptx"'}
    )

def _export_pdf(title: str, content: dict) -> StreamingResponse:
    from fpdf import FPDF
    import io
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=10)
    
    # Title
    safe_title = title.encode('latin-1', 'replace').decode('latin-1')
    pdf.set_font("helvetica", style="B", size=14)
    pdf.cell(0, 10, safe_title, new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("helvetica", size=10)
    
    cells = content.get("cells", {})
    if cells:
        max_row = max((int(r) for r in cells.keys()), default=-1)
        max_col = max((int(c) for r in cells.values() for c in r.keys()), default=-1)
        
        # Calculate optimal column widths (simplified)
        col_widths = [20] * (max_col + 1)
        for c in range(max_col + 1):
            max_len = 5 # min width
            for r in range(max_row + 1):
                val = str(cells.get(str(r), {}).get(str(c), ""))
                if len(val) > max_len:
                    max_len = len(val)
            col_widths[c] = min(max_len * 2.5, 60) 

        for r in range(max_row + 1):
            for c in range(max_col + 1):
                val = str(cells.get(str(r), {}).get(str(c), ""))
                safe_val = val.encode('latin-1', 'replace').decode('latin-1')
                pdf.cell(col_widths[c], 8, safe_val, border=1)
            pdf.ln(8)

    buf = io.BytesIO(pdf.output())
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{title}.pdf"'}
    )